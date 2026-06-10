import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Colors & Style Constants
COLOR_TEAL = RGBColor(0, 128, 128)        # Accent/Header Teal
COLOR_DARK_SLATE = RGBColor(46, 59, 78)  # Body text Dark Slate
COLOR_PALE_MINT = RGBColor(230, 245, 240) # Center Panel Background & Alt Row Mint
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GREY = RGBColor(120, 120, 120)

FONT_FAMILY = "Arial"

# 2. Poster Contents
TITLE_TEXT = "Behavioral Diversity in Mobile Money Ledgers Outpredicts Transaction Sequences for National-Scale Credit Risk"
AUTHORS_TEXT = "Attabra Benjamin Ekow (KNUST, YARA Fellow) — presenting author (beattabra@st.knust.edu.gh)\nCo-authors: [Name 1], [Name 2]"
AFFILIATIONS_TEXT = "Kwame Nkrumah University of Science and Technology (KNUST), YARA Fellows Program, Google AI Community Centre (Accra, Ghana)"

# Column 1 Contents (Introduction, Research Question, Methods)
INTRO_HEADER = "1. INTRODUCTION / MOTIVATION"
INTRO_BODY = [
    "Traditional credit bureau infrastructure in sub-Saharan Africa remains thin or absent for the vast majority of the population, limiting financial inclusion.",
    "Mobile money (MoMo) transactional logs provide a rich, alternative data source for credit risk profiling.",
    "Using a two-stage strategy, we first built a synthetic MoMo benchmark (Paper A) to design a leakage-free temporal pipeline.",
    "In this study (Paper B), we scale and validate our framework on a national-scale database of 374 million transaction logs from Telecel Ghana to determine if sequential deep learning holds any predictive advantage over static, well-engineered behavioral aggregates."
]

RESEARCH_HEADER = "2. RESEARCH QUESTION & HYPOTHESIS"
RESEARCH_BODY = [
    "Core Research Question:\nDo sequential deep learning models (LSTMs/GRUs) outperform static aggregate classifiers when predicting mobile money credit default at national scale?",
    "Hypothesis:\nRecipient entropy (behavioral diversity) dominates transaction-level sequence order as a predictor of default risk."
]

METHODS_HEADER = "3. METHODS & LEAKAGE-FREE PIPELINE"
METHODS_BODY = [
    "We processed 374,295,424 transaction logs covering 474,312 borrowers across a 90-day window (Sep – Nov 2025) on Azure Databricks.",
    "For each borrower, the index loan is set as their last loan disbursement (last_loan_ts). Features are built strictly from transactions before the cutoff; labels are derived from events after the cutoff."
]

# Center Panel Contents (Hero Main Finding)
HERO_FINDING_HEADER = "THE CENTRAL FINDING:"
HERO_FINDING_TEXT = "Borrower credit risk is written in the diversity of their transaction network (recipient entropy), NOT in the sequential patterns of their transactions."

RESULTS_TABLE_DATA = [
    ["Model", "y_default AUC-ROC", "y_default AUC-PR", "y_default ECE", "y_bad AUC-ROC", "y_bad AUC-PR"],
    ["LightGBM", "0.9464", "0.9462", "0.0058", "0.8069", "0.9291"],
    ["XGBoost", "0.9462", "0.9460", "0.0060", "0.8069", "0.9289"],
    ["Random Forest", "0.9370", "0.9346", "0.0568", "0.7905", "0.9223"],
    ["Logistic Regression", "0.9290", "0.9246", "0.0830", "0.6996", "0.8732"],
    ["GRU (Sequence)", "0.7553", "0.7120", "0.1240", "0.6510", "0.8122"]
]

CENTER_DETAIL_TEXT = [
    "Probability Calibration & XAI:",
    "LightGBM and XGBoost achieved near-perfect calibration (ECE < 0.006) which is vital for real-world automated loan pricing, whereas sequential networks showed severe calibration degradation (ECE > 0.12).",
    "Feature ablation demonstrates that dropping the Behavioral Diversity group causes a significant -0.035 drop in AUC-ROC, confirming it as the most powerful non-loan risk signal."
]

# Column 4 Contents (Conclusions, References, Acknowledgements)
CONCLUSIONS_HEADER = "5. CONCLUSIONS"
CONCLUSIONS_BODY = [
    "• Static beats Sequence: Sequential deep learning models (LSTMs/GRUs) add no discriminative or calibration value over well-engineered aggregate statistics on mobile money data.",
    "• Behavioral Diversity is the Key Signal: The entropy and diversity of a user's transaction network (e.g., number of unique transfer recipients) serves as a robust proxy for economic stability and financial inclusion depth.",
    "• Production Calibration: Gradient-boosted trees (LightGBM/XGBoost) are exceptionally well-calibrated, making them the optimal choices for real-world automated credit scoring.",
    "• Limitations & Future Work: Documented 463 cold-start borrowers with zero history; future work will test if behavioral entropy can act as a \"proxy history\" for zero-history users."
]

REFERENCES_HEADER = "6. KEY REFERENCES"
REFERENCES_BODY = [
    "[1] Bank of Ghana. (2025). Summary of Economic and Financial Data.",
    "[2] GSMA. (2025). The State of the Industry Report on Mobile Money.",
    "[3] Morrison, M. (2020). The Better Poster: Layouts for Academic Conferences.",
    "[4] Attabra, B. E. (2026). Sequential Deep Learning for Credit Risk Modeling in Data-Constrained Environments (Paper A)."
]

ACKNOWLEDGEMENTS_HEADER = "7. ACKNOWLEDGEMENTS & CONTACT"
ACKNOWLEDGEMENTS_BODY = [
    "Supported by the YARA Fellows Program. Computing resources provided by Azure Databricks. Special thanks to Telecel Ghana for providing anonymized mobile money logs for research purposes.",
    "Contact: beattabra@st.knust.edu.gh | https://github.com/Attabeezy/seqcredit-model"
]

def set_shape_text(shape, text, font_size=24, bold=False, color=COLOR_DARK_SLATE, italic=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_FAMILY
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color

def main():
    pptx_template = r"C:\Users\Dell\Projects\seqcredit_model\docs\publications\yara\PosterPresentations.com-36x48-Template-Marquee.pptx"
    output_path = r"C:\Users\Dell\Projects\seqcredit_model\docs\publications\yara\YARA_Poster_36x48.pptx"

    if not os.path.exists(pptx_template):
        print(f"Error: Template not found at {pptx_template}")
        return

    prs = Presentation(pptx_template)
    slide = prs.slides[0]

    # Map layout idx to slide shapes
    shapes_by_idx = {}
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_by_idx[shape.placeholder_format.idx] = shape

    # 3. Update Title, Authors, and Affiliations (Title = 72pt, Authors = 28pt, Affils = 24pt)
    if 153 in shapes_by_idx:
        print("Writing Title...")
        set_shape_text(shapes_by_idx[153], TITLE_TEXT, font_size=72, bold=True, color=COLOR_TEAL, align=PP_ALIGN.CENTER)
        # Shift Title up slightly and make it taller to prevent wrapping issues
        shapes_by_idx[153].top = Inches(0.2)
        shapes_by_idx[153].height = Inches(1.8)
    if 151 in shapes_by_idx:
        print("Writing Authors...")
        set_shape_text(shapes_by_idx[151], AUTHORS_TEXT, font_size=28, bold=True, color=COLOR_DARK_SLATE, align=PP_ALIGN.CENTER)
        shapes_by_idx[151].top = Inches(2.1)
        shapes_by_idx[151].height = Inches(1.0)
    if 150 in shapes_by_idx:
        print("Writing Affiliations...")
        set_shape_text(shapes_by_idx[150], AFFILIATIONS_TEXT, font_size=24, bold=False, color=COLOR_GREY, italic=True, align=PP_ALIGN.CENTER)
        shapes_by_idx[150].top = Inches(3.2)
        shapes_by_idx[150].height = Inches(0.8)

    # 4. Hide / Clear Middle Column Placeholders (Columns 2 & 3: idx 22, 21, 24, 23)
    # This leaves the center area empty for our custom design.
    for m_idx in [22, 21, 24, 23]:
        if m_idx in shapes_by_idx:
            sh = shapes_by_idx[m_idx]
            sh.text_frame.clear()
            sh.left = Inches(-100) # Move off-slide

    # 5. Position & Populate Column 1 (Left Column, X = 1.3 inches, Width = 11.0 inches)
    # We will lay out Intro, Research Q, and Methods in Column 1.
    print("Populating Column 1...")
    
    # 5.1 Introduction Header
    if 11 in shapes_by_idx:
        header_sh = shapes_by_idx[11]
        header_sh.left = Inches(1.3)
        header_sh.width = Inches(11.0)
        header_sh.top = Inches(5.0)
        header_sh.height = Inches(0.8)
        set_shape_text(header_sh, INTRO_HEADER, font_size=36, bold=True, color=COLOR_TEAL)
        
    # 5.2 Introduction Body
    if 10 in shapes_by_idx:
        body_sh = shapes_by_idx[10]
        body_sh.left = Inches(1.3)
        body_sh.width = Inches(11.0)
        body_sh.top = Inches(5.8)
        body_sh.height = Inches(5.2)
        tf = body_sh.text_frame
        tf.clear()
        for idx, para in enumerate(INTRO_BODY):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.name = FONT_FAMILY
            p.font.size = Pt(24)
            p.font.color.rgb = COLOR_DARK_SLATE
            p.space_after = Pt(12)

    # 5.3 Research Question Header
    if 20 in shapes_by_idx:
        header_sh = shapes_by_idx[20]
        header_sh.left = Inches(1.3)
        header_sh.width = Inches(11.0)
        header_sh.top = Inches(11.5)
        header_sh.height = Inches(0.8)
        set_shape_text(header_sh, RESEARCH_HEADER, font_size=36, bold=True, color=COLOR_TEAL)
        
    # 5.4 Research Question Body
    if 96 in shapes_by_idx:
        body_sh = shapes_by_idx[96]
        body_sh.left = Inches(1.3)
        body_sh.width = Inches(11.0)
        body_sh.top = Inches(12.3)
        body_sh.height = Inches(3.2)
        tf = body_sh.text_frame
        tf.clear()
        for idx, para in enumerate(RESEARCH_BODY):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.name = FONT_FAMILY
            p.font.size = Pt(24)
            p.font.color.rgb = COLOR_DARK_SLATE
            p.space_after = Pt(12)

    # 5.5 Methods Header (dynamically added text box)
    methods_hdr_box = slide.shapes.add_textbox(Inches(1.3), Inches(16.0), Inches(11.0), Inches(0.8))
    set_shape_text(methods_hdr_box, METHODS_HEADER, font_size=36, bold=True, color=COLOR_TEAL)

    # 5.6 Methods Body (dynamically added text box)
    methods_body_box = slide.shapes.add_textbox(Inches(1.3), Inches(16.8), Inches(11.0), Inches(8.5))
    tf_m = methods_body_box.text_frame
    tf_m.word_wrap = True
    tf_m.clear()
    for idx, para in enumerate(METHODS_BODY):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        p.text = para
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_DARK_SLATE
        p.space_after = Pt(12)

    # 5.7 Visual Flowchart (Timeline Partitioning Diagram in Column 1, Y = 26.0 to 32.0)
    print("Drawing Methods Flowchart...")
    diagram_top = 25.8
    box_width = Inches(3.0)
    box_height = Inches(1.8)
    y_offset = Inches(diagram_top)

    # Box 1: Pre-Loan (Features)
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), y_offset, box_width, box_height)
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_PALE_MINT
    box1.line.color.rgb = COLOR_TEAL
    box1.line.width = Pt(2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.clear()
    p1 = tf1.paragraphs[0]
    p1.text = "1. PRE-LOAN LOGS\n(Features)\n\n38 Aggregates\n19 Sequence Features"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_DARK_SLATE

    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), y_offset + Inches(0.6), Inches(0.4), Inches(0.6))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_TEAL
    arrow1.line.fill.background()

    # Box 2: Index Loan Cutoff (Teal)
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), y_offset, box_width, box_height)
    box2.fill.solid()
    box2.fill.fore_color.rgb = COLOR_TEAL
    box2.line.fill.background()
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "2. INDEX LOAN\n(Cutoff Point)\n\nStrict leakage-free\ntime split at last_loan_ts"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.3), y_offset + Inches(0.6), Inches(0.4), Inches(0.6))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_TEAL
    arrow2.line.fill.background()

    # Box 3: Post-Loan (Labels)
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), y_offset, box_width, box_height)
    box3.fill.solid()
    box3.fill.fore_color.rgb = COLOR_PALE_MINT
    box3.line.color.rgb = COLOR_TEAL
    box3.line.width = Pt(2)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    tf3.clear()
    p3 = tf3.paragraphs[0]
    p3.text = "3. POST-LOAN\n(Labels)\n\nRepayment history\ny_default / y_bad"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_DARK_SLATE


    # 6. Build Center Panel (Hero Zone, X = 13.0 to 35.0, Width = 22.0)
    print("Constructing Center Panel (Better Poster format)...")
    
    # 6.1 Panel Background Shape (Rounded Rectangle)
    panel_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(13.0), Inches(5.0), Inches(22.0), Inches(28.0))
    panel_bg.fill.solid()
    panel_bg.fill.fore_color.rgb = COLOR_PALE_MINT
    panel_bg.line.color.rgb = COLOR_TEAL
    panel_bg.line.width = Pt(3)

    # 6.2 Giant Center Main Finding
    finding_hdr_box = slide.shapes.add_textbox(Inches(13.5), Inches(5.3), Inches(21.0), Inches(0.6))
    set_shape_text(finding_hdr_box, HERO_FINDING_HEADER, font_size=28, bold=True, color=COLOR_TEAL, align=PP_ALIGN.CENTER)

    finding_body_box = slide.shapes.add_textbox(Inches(13.5), Inches(6.0), Inches(21.0), Inches(3.0))
    # Giant text size = 44pt for the central finding billboard
    set_shape_text(finding_body_box, HERO_FINDING_TEXT, font_size=44, bold=True, color=COLOR_DARK_SLATE, align=PP_ALIGN.CENTER)

    # 6.3 Results Table in Center Panel
    print("Adding Results Table in Center Panel...")
    table_left = Inches(13.5)
    table_top = Inches(9.4)
    table_width = Inches(21.0)
    table_height = Inches(4.2)

    table_shape = slide.shapes.add_table(6, 6, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Adjust Column Widths
    table.columns[0].width = Inches(5.0)  # Model name
    for c in range(1, 6):
        table.columns[c].width = Inches(3.2)

    # Populate Table
    for r_idx, row_data in enumerate(RESULTS_TABLE_DATA):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            p.font.name = FONT_FAMILY
            p.font.size = Pt(20) # High legibility size for table cells
            
            # Apply Styles
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TEAL
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
            else:
                cell.fill.solid()
                if r_idx % 2 == 1:
                    cell.fill.fore_color.rgb = COLOR_WHITE
                else:
                    cell.fill.fore_color.rgb = COLOR_PALE_MINT
                p.font.color.rgb = COLOR_DARK_SLATE
                if r_idx == 1: # Highlight LightGBM (our SOTA baseline)
                    p.font.bold = True

    # 6.4 Visual Feature Importance Bar Chart (Y = 14.2 to 20.0)
    print("Drawing Horizontal Bar Chart for Feature Ablation...")
    chart_title_box = slide.shapes.add_textbox(Inches(13.5), Inches(14.0), Inches(21.0), Inches(0.8))
    set_shape_text(chart_title_box, "FEATURE ABLATION SENSITIVITY (AUC-ROC drop when feature group is removed)", font_size=22, bold=True, color=COLOR_TEAL, align=PP_ALIGN.CENTER)

    features_ablation = [
        ("Behavioral Diversity (recipient entropy, type count)", 0.035, COLOR_TEAL),
        ("Loan History (previous default/payment volume)", 0.020, COLOR_GREY),
        ("Activity Intensity (daily transaction frequency)", 0.005, COLOR_GREY),
        ("All Other Groups (amounts, fees, timing)", 0.001, COLOR_GREY)
    ]

    for idx, (feat_name, delta, color) in enumerate(features_ablation):
        y_pos = 15.0 + idx * 1.1
        
        # Label Text
        lbl_box = slide.shapes.add_textbox(Inches(13.5), Inches(y_pos), Inches(9.0), Inches(0.7))
        set_shape_text(lbl_box, feat_name, font_size=18, bold=True, color=COLOR_DARK_SLATE)
        
        # Draw Visual Bar
        # Max scale: delta = 0.035 matches 9.5 inches width
        bar_width = Inches((delta / 0.035) * 9.5)
        bar_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(22.5), Inches(y_pos + 0.1), bar_width, Inches(0.5))
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = color
        bar_shape.line.fill.background()
        
        # Value Label
        val_box = slide.shapes.add_textbox(Inches(22.6) + bar_width, Inches(y_pos), Inches(1.5), Inches(0.7))
        set_shape_text(val_box, f"-{delta:.3f}", font_size=18, bold=True, color=color)

    # 6.5 Bottom Row: Finding details & QR Code side-by-side (Y = 20.0 to 32.5)
    print("Constructing Bottom Row inside Center Panel...")
    
    # Left Block (Details & Calibration Discussion)
    details_box = slide.shapes.add_textbox(Inches(13.5), Inches(20.0), Inches(13.0), Inches(11.5))
    tf_d = details_box.text_frame
    tf_d.word_wrap = True
    tf_d.clear()
    for idx, para in enumerate(CENTER_DETAIL_TEXT):
        p = tf_d.paragraphs[0] if idx == 0 else tf_d.add_paragraph()
        p.text = para
        p.font.name = FONT_FAMILY
        if "Probability Calibration" in para:
            p.font.bold = True
            p.font.size = Pt(22)
            p.font.color.rgb = COLOR_TEAL
            p.space_before = Pt(12)
        else:
            p.font.size = Pt(24)
            p.font.color.rgb = COLOR_DARK_SLATE
        p.space_after = Pt(12)

    # Right Block: QR Code Box Visual
    qr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(28.0), Inches(20.2), Inches(6.0), Inches(5.8))
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = COLOR_WHITE
    qr_box.line.color.rgb = COLOR_TEAL
    qr_box.line.width = Pt(2)
    tf_qr = qr_box.text_frame
    tf_qr.word_wrap = True
    tf_qr.clear()
    p_qr = tf_qr.paragraphs[0]
    p_qr.text = "[ QR CODE PLACEHOLDER ]\n\nScan here to access:\n• Paper B preprint\n• Python source package\n• Databricks pipeline configs\n• Calibration curves"
    p_qr.alignment = PP_ALIGN.CENTER
    p_qr.font.name = FONT_FAMILY
    p_qr.font.size = Pt(16)
    p_qr.font.bold = True
    p_qr.font.color.rgb = COLOR_DARK_SLATE

    # Right Block: Partner Logo Box Visual
    logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(28.0), Inches(26.6), Inches(6.0), Inches(5.4))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLOR_DARK_SLATE
    logo_box.line.fill.background()
    tf_l = logo_box.text_frame
    tf_l.word_wrap = True
    tf_l.clear()
    p_l = tf_l.paragraphs[0]
    p_l.text = "TELECEL GHANA\n\n[ Data Partner Badge ]\n\n374M Transaction Logs\nNational Scale Validation"
    p_l.alignment = PP_ALIGN.CENTER
    p_l.font.name = FONT_FAMILY
    p_l.font.size = Pt(16)
    p_l.font.bold = True
    p_l.font.color.rgb = COLOR_WHITE


    # 7. Position & Populate Column 4 (Right Column, X = 35.8 inches, Width = 11.0 inches)
    # We will lay out Conclusions, References, and Acknowledgements/Contact here.
    print("Populating Column 4...")
    
    # 7.1 Conclusions Header
    if 25 in shapes_by_idx:
        hdr_sh = shapes_by_idx[25]
        hdr_sh.left = Inches(35.8)
        hdr_sh.width = Inches(11.0)
        hdr_sh.top = Inches(5.0)
        hdr_sh.height = Inches(0.8)
        set_shape_text(hdr_sh, CONCLUSIONS_HEADER, font_size=36, bold=True, color=COLOR_TEAL)
        
    # 7.2 Conclusions Body
    if 26 in shapes_by_idx:
        body_sh = shapes_by_idx[26]
        body_sh.left = Inches(35.8)
        body_sh.width = Inches(11.0)
        body_sh.top = Inches(5.8)
        body_sh.height = Inches(9.2)
        tf = body_sh.text_frame
        tf.clear()
        for idx, para in enumerate(CONCLUSIONS_BODY):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.name = FONT_FAMILY
            p.font.size = Pt(24)
            p.font.color.rgb = COLOR_DARK_SLATE
            p.space_after = Pt(14)

    # 7.3 References Header
    if 27 in shapes_by_idx:
        hdr_sh = shapes_by_idx[27]
        hdr_sh.left = Inches(35.8)
        hdr_sh.width = Inches(11.0)
        hdr_sh.top = Inches(15.5)
        hdr_sh.height = Inches(0.8)
        set_shape_text(hdr_sh, REFERENCES_HEADER, font_size=36, bold=True, color=COLOR_TEAL)
        
    # 7.4 References Body (smaller size for space efficiency)
    if 28 in shapes_by_idx:
        body_sh = shapes_by_idx[28]
        body_sh.left = Inches(35.8)
        body_sh.width = Inches(11.0)
        body_sh.top = Inches(16.3)
        body_sh.height = Inches(7.7)
        tf = body_sh.text_frame
        tf.clear()
        for idx, para in enumerate(REFERENCES_BODY):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14) # References at 14pt (readable but space-saving)
            p.font.color.rgb = COLOR_DARK_SLATE
            p.space_after = Pt(8)

    # 7.5 Acknowledgements & Contact Header
    if 29 in shapes_by_idx:
        hdr_sh = shapes_by_idx[29]
        hdr_sh.left = Inches(35.8)
        hdr_sh.width = Inches(11.0)
        hdr_sh.top = Inches(24.5)
        hdr_sh.height = Inches(0.8)
        set_shape_text(hdr_sh, ACKNOWLEDGEMENTS_HEADER, font_size=36, bold=True, color=COLOR_TEAL)
        
    # 7.6 Acknowledgements & Contact Body
    if 30 in shapes_by_idx:
        body_sh = shapes_by_idx[30]
        body_sh.left = Inches(35.8)
        body_sh.width = Inches(11.0)
        body_sh.top = Inches(25.3)
        body_sh.height = Inches(7.7)
        tf = body_sh.text_frame
        tf.clear()
        for idx, para in enumerate(ACKNOWLEDGEMENTS_BODY):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.name = FONT_FAMILY
            p.font.size = Pt(18) # Acknowledgements at 18pt
            p.font.color.rgb = COLOR_DARK_SLATE
            p.space_after = Pt(10)

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated and saved Better Poster to {output_path}")

if __name__ == "__main__":
    main()
